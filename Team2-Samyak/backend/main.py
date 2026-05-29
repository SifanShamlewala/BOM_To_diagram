"""
DeltaV Architecture Editor - Using Friend's Pipeline
BOM → archEngine → layoutEngine → pptRenderer
"""

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Dict, Any
import io
import pandas as pd

# Import friend's pipeline (CORRECT SPELLING)
from archEngine import ArchitectureEngine, ArchitectureConfig
from layoutEngine import LayoutEngine
from pptRenderer import PptRenderer, RenderConfig

app = FastAPI(title="DeltaV Editor with Friend's Pipeline", version="4.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============================================================================
# BOM Parser → Semantic JSON (All 6 Rules)
# ============================================================================

def parse_bom_to_semantic_json(file_bytes: bytes) -> dict:
    """
    Parse BOM Excel to Semantic JSON format
    Following all 6 corrected rules
    """
    df = pd.read_excel(io.BytesIO(file_bytes))
    df = df.dropna(subset=['Sr. No'])
    
    # Rule 1: Determine rooms dynamically
    rooms = ["OTHER"]
    
    operator_items = df[
        (df['Category'].isin(['MON', 'PRT', 'CON'])) & 
        (~df['Description'].str.contains('DTS', na=False, case=False))
    ]
    tower_ws = df[
        (df['Category'] == 'WS') & 
        (df['Description'].str.contains('Tower', na=False, case=False))
    ]
    
    if len(operator_items) > 0 or len(tower_ws) > 0:
        rooms.append("OPERATOR")
    
    # Build semantic JSON
    semantic_rooms = []
    
    for room_type in rooms:
        room_data = {
            "room_type": room_type,
            "CNTR": get_controllers(df, room_type),
            "CHARM": get_charm(df, room_type),
            "FOPP": get_fopp(df, room_type),
            "ENC_COMP": get_enc_comp(df, room_type)
        }
        semantic_rooms.append(room_data)
    
    return {"rooms": semantic_rooms}

def get_controllers(df, room_type):
    """Rule 2: CNTR - quantity=1 if any CNTR items; redundant=true if 'Redundant' in description"""
    if room_type == "OPERATOR":
        return {"quantity": 0, "redundant": False}
    
    cntr_items = df[df['Category'] == 'CNTR']
    
    if len(cntr_items) == 0:
        return {"quantity": 0, "redundant": False}
    
    redundant = any('Redundant' in str(desc) for desc in cntr_items['Description'])
    
    return {
        "quantity": 1,
        "redundant": redundant
    }

def get_charm(df, room_type):
    """Rule 3: CHARM - ONLY 'Baseplate Assembly' item → baseplates count"""
    if room_type == "OPERATOR":
        return {"baseplates": 0}
    
    charm_items = df[df['Category'] == 'CHARM']
    baseplate_item = charm_items[
        charm_items['Description'].str.contains('Baseplate Assembly', na=False, case=False)
    ]
    
    if baseplate_item.empty:
        return {"baseplates": 0}
    
    return {"baseplates": int(baseplate_item.iloc[0]['Qty'])}

def get_fopp(df, room_type):
    """Rule 4: FOPP - Include BOTH Category=FOPP AND Category=MC; count unique TYPES"""
    fopp_mc_items = df[df['Category'].isin(['FOPP', 'MC'])]
    quantity = len(fopp_mc_items) if room_type == "OTHER" else 0
    return {"quantity": quantity}

def get_enc_comp(df, room_type):
    """Rule 5: ENC_COMP - Categories WS, SWT, SCR, UPS, MON, PRT; simplify names; actual BOM qty"""
    enc_comp_categories = ['WS', 'SWT', 'SCR', 'UPS', 'MON', 'PRT']
    items = df[df['Category'].isin(enc_comp_categories)]
    
    enc_comp = []
    
    for _, row in items.iterrows():
        # Rule 6: Exclude DTS items
        if 'DTS' in str(row['Description']).upper():
            continue
        
        part_name, belongs = classify_enc_comp(row['Category'], row['Description'], room_type)
        
        if belongs:
            enc_comp.append({
                "part_name": part_name,
                "quantity": int(row['Qty'])
            })
    
    return enc_comp

def classify_enc_comp(category: str, description: str, room_type: str):
    """Classify ENC_COMP item to room with simplified names"""
    desc_lower = description.lower()
    
    if room_type == "OTHER":
        if category == 'WS' and 'rack' in desc_lower:
            return ("DeltaV Rack Workstation", True)
        elif category == 'SWT':
            if 'deltav' in desc_lower:
                return ("DeltaV Smart Switch", True)
            elif 'hirschmann' in desc_lower:
                return ("Hirschmann Industrial Switch", True)
        elif category == 'SCR':
            return ("19 inch Sliding Screen", True)
        elif category == 'UPS':
            return ("APC UPS", True)
    
    elif room_type == "OPERATOR":
        if category == 'WS' and 'tower' in desc_lower:
            return ("DeltaV Full-sized Tower Workstation", True)
        elif category == 'MON' and '24' in description:
            return ("24-inch Monitor", True)
        elif category == 'PRT':
            return ("Printer", True)
    
    return ("", False)

# ============================================================================
# API Endpoints
# ============================================================================

@app.get("/")
def root():
    return {
        "status": "online",
        "service": "DeltaV Editor with Friend's Pipeline",
        "version": "4.0",
        "pipeline": "BOM → archEngine → layoutEngine → pptRenderer"
    }

@app.post("/api/parse-bom")
async def parse_bom_endpoint(file: UploadFile = File(...)):
    """
    Upload BOM → Get structured data for web editor
    Uses friend's archEngine to build architecture
    """
    try:
        contents = await file.read()
        print(f"📥 Received BOM: {file.filename}")
        
        # Step 1: Parse to semantic JSON (our 6 rules)
        semantic = parse_bom_to_semantic_json(contents)
        print(f"📋 Semantic JSON: {len(semantic['rooms'])} rooms")
        
        # Step 2: Use friend's archEngine to build architecture
        # Step 2: Use friend's archEngine to build architecture
        arch_engine = ArchitectureEngine(ArchitectureConfig())
        architecture = arch_engine.build(semantic)
        print(f"🏗️  Architecture built: {len(architecture['rooms'])} rooms")

        # DEBUG: Print actual structure
        import json
        print("🔍 DEBUG - Architecture structure:")
        print(json.dumps(architecture, indent=2, default=str))
        
        # Convert to simplified format for frontend
        simplified_data = {
            "rooms": []
        }
        
        for room in architecture['rooms']:
            room_data = {
                "room_type": room['room_type'],
                "controllers": [],
                "baseplates": 0,
                "fopp_count": 0,
                "devices": []
            }
            
            # Extract cabinets
            # Extract cabinets
            for cabinet in room.get('cabinets', []):
                # Controllers
                if cabinet['type'] == 'SERVER':
                    for device in cabinet.get('devices', []):
                        # Safe access with .get() instead of ['key']
                        device_id = device.get('node_id') or device.get('device_id') or device.get('id', 'DEVICE')
                        device_name = device.get('part_name') or device.get('name', 'Unknown Device')
                        
                        if 'CNTR' in device_id or 'controller' in device_name.lower():
                            room_data['controllers'].append({
                                "id": device_id,
                                "label": device_id,
                                "type": "controller"
                            })
                        else:
                            room_data['devices'].append({
                                "id": device_id,
                                "label": device_name,
                                "type": "device",
                                "qty": 1
                            })
                
                # Baseplates from I/O cabinets
                elif cabinet['type'] == 'IO':
                    for tower in cabinet.get('towers', []):
                        room_data['baseplates'] += len(tower.get('baseplates', []))
                
                # Workdesk devices
                elif cabinet['type'] == 'WORKDESK':
                    for device in cabinet.get('devices', []):
                        device_id = device.get('node_id') or device.get('device_id') or device.get('id', 'DEVICE')
                        device_name = device.get('part_name') or device.get('name', 'Unknown Device')
                        
                        room_data['devices'].append({
                            "id": device_id,
                            "label": device_name,
                            "type": "device",
                            "qty": 1
                        })
            
            # FOPP
            room_data['fopp_count'] = len(room.get('fopp_nodes', []))
            
            simplified_data['rooms'].append(room_data)
        
        return {"status": "success", "data": simplified_data}
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

class ExportRequest(BaseModel):
    shapes: List[Dict[str, Any]]
    connections: List[Dict[str, str]] = []

@app.post("/api/export-ppt")
async def export_ppt_from_positions(request: ExportRequest):
    """Export PPT using browser positions - reflects all added components"""
    try:
        shapes = request.shapes
        print(f"📊 Exporting PPT with {len(shapes)} shapes")

        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN as TEXT_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        from lxml import etree

        prs = Presentation()
        prs.slide_width = Inches(20)   # Wide canvas to fit large diagrams
        prs.slide_height = Inches(12)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(6), Inches(0.5))
        tf = txBox.text_frame
        tf.text = "DeltaV Architecture Diagram"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(30, 30, 30)

        # Color palette
        colors = {
            "room":        (230, 240, 255),  # Light blue-grey
            "cabinet":     (70,  120, 200),  # Medium blue
            "controller":  (100, 180, 240),  # Sky blue
            "device":      (210, 225, 250),  # Very light blue
            "baseplate":   (160, 220, 170),  # Light green
            "fopp":        (255, 160,  80),  # Orange
        }
        border_colors = {
            "room":       (120, 150, 200),
            "cabinet":    (40,  80,  160),
            "controller": (50, 130, 190),
            "device":     (130, 150, 200),
            "baseplate":  (80,  160,  90),
            "fopp":       (200, 100,  30),
        }

        # Scale: browser px → PPT inches (browser ~96dpi, scale to fit)
        # Find bounding box of all shapes
        if shapes:
            max_x = max(s['x'] + s.get('width', 100) for s in shapes)
            max_y = max(s['y'] + s.get('height', 50) for s in shapes)
        else:
            max_x, max_y = 1920, 1080

        canvas_w = float(prs.slide_width) / 914400  # EMU → inches
        canvas_h = float(prs.slide_height) / 914400
        margin = 0.5
        usable_w = canvas_w - margin * 2
        usable_h = canvas_h - margin * 2 - 0.7  # leave room for title

        scale = min(usable_w / (max_x / 96), usable_h / (max_y / 96))

        def px_to_inches(px):
            return px / 96 * scale

        # Sort: rooms first, then cabinets, then contents (so rooms are drawn behind)
        type_order = {"room": 0, "cabinet": 1, "fopp": 2, "controller": 3, "baseplate": 4, "device": 5}
        sorted_shapes = sorted(shapes, key=lambda s: type_order.get(s['type'], 6))

        for shape in sorted_shapes:
            left   = Inches(margin + px_to_inches(shape['x']))
            top    = Inches(margin + 0.7 + px_to_inches(shape['y']))
            width  = Inches(max(px_to_inches(shape.get('width', 100)), 0.15))
            height = Inches(max(px_to_inches(shape.get('height', 40)), 0.1))

            fill_rgb  = RGBColor(*colors.get(shape['type'], (200, 200, 200)))
            line_rgb  = RGBColor(*border_colors.get(shape['type'], (80, 80, 80)))
            line_w    = Pt(1.5 if shape['type'] == 'room' else 0.75)

            from pptx.enum.shapes import MSO_SHAPE
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = fill_rgb
            rect.line.color.rgb = line_rgb
            rect.line.width = line_w

            # Transparency for rooms so children are visible
            if shape['type'] == 'room':
                # Set fill transparency via XML
                sp = rect._element
                sp_pr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                if sp_pr is not None:
                    srgb = sp_pr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    if srgb is None:
                        srgb = etree.SubElement(sp_pr, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    alpha.set('val', '50000')  # 50% opacity

            tf = rect.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = shape['label']

            font_size = {
                "room": 11, "cabinet": 9, "controller": 7,
                "device": 7, "baseplate": 6, "fopp": 8
            }.get(shape['type'], 7)

            p.font.size = Pt(font_size)
            p.font.bold = shape['type'] in ['room', 'cabinet']
            p.font.color.rgb = RGBColor(255, 255, 255) if shape['type'] == 'cabinet' else RGBColor(20, 20, 20)
            p.alignment = TEXT_ALIGN.CENTER

            from pptx.util import Pt as PtUtil
            tf.margin_top = PtUtil(2)
            tf.margin_bottom = PtUtil(2)
            tf.margin_left = PtUtil(3)
            tf.margin_right = PtUtil(3)

        # Add legend
        legend_x = Inches(canvas_w - 2.5)
        legend_y = Inches(0.15)
        legend_items = [("Room", colors["room"]), ("Cabinet", colors["cabinet"]),
                        ("Controller", colors["controller"]), ("Device", colors["device"]),
                        ("Baseplate", colors["baseplate"]), ("FOPP", colors["fopp"])]
        for i, (label, color) in enumerate(legend_items):
            bx = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                legend_x, legend_y + Inches(i * 0.28), Inches(0.22), Inches(0.2))
            bx.fill.solid()
            bx.fill.fore_color.rgb = RGBColor(*color)
            bx.line.color.rgb = RGBColor(80, 80, 80)
            bx.line.width = Pt(0.5)
            tx = slide.shapes.add_textbox(
                legend_x + Inches(0.28), legend_y + Inches(i * 0.28), Inches(1.5), Inches(0.2))
            tx.text_frame.text = label
            tx.text_frame.paragraphs[0].font.size = Pt(8)

        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)

        return StreamingResponse(
            ppt_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=deltav_architecture.pptx"}
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    print("\n" + "="*70)
    print("  DeltaV Editor - Using Friend's Pipeline")
    print("  BOM → archEngine → layoutEngine → pptRenderer")
    print("="*70 + "\n")
    uvicorn.run(app, host="0.0.0.0", port=8000)
