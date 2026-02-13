from pptx import Presentation
from pptx.util import Inches

class PptRenderer:
    def __init__(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

    def draw_room(self, room):
        self.slide.shapes.add_shape(
            autoshape_type_id=1,  # rectangle
            left=Inches(room["x"]/96),
            top=Inches(room["y"]/96),
            width=Inches(room["width"]/96),
            height=Inches(room["height"]/96)
        )

    def draw_enclosure(self, encl):
        self.slide.shapes.add_shape(
            autoshape_type_id=1,
            left=Inches(encl["x"]/96),
            top=Inches(encl["y"]/96),
            width=Inches(encl["width"]/96),
            height=Inches(encl["height"]/96)
        )

    def draw_device(self, dev):
        # Replace with icon insert later
        self.slide.shapes.add_shape(
            autoshape_type_id=1,
            left=Inches(dev["x"]/96),
            top=Inches(dev["y"]/96),
            width=Inches(dev.get("width", 40)/96),
            height=Inches(dev.get("height", 20)/96)
        )

    def draw_link(self, src, dst):
        self.slide.shapes.add_connector(
            1,
            Inches(src["x"]/96),
            Inches(src["y"]/96),
            Inches(dst["x"]/96),
            Inches(dst["y"]/96)
        )

    def save(self, path):
        self.prs.save(path)
