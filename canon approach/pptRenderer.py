"""
DeltaV DCS PPT Renderer
Stage 4: Layout JSON → PowerPoint diagram

Responsibilities:
- Convert row/col → x/y coordinates (in inches)
- Draw room boundaries, cabinets, towers, devices, FOPP nodes
- Draw connections (FOPP_TO_CABINET, INTER_ROOM_FIBER)
- Zero engineering logic — pure visual translation

Uses python-pptx. All geometry constants live in RenderConfig.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from dataclasses import dataclass
from typing import Dict, Tuple
import copy


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def rgb(hex_str: str) -> RGBColor:
    """Convert 6-char hex string to RGBColor."""
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def add_rect(slide, x, y, w, h, fill_hex, border_hex, border_width_pt=1.0):
    """Add a filled rectangle with border."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.color.rgb = rgb(border_hex)
    shape.line.width = Pt(border_width_pt)
    return shape


def add_label(slide, x, y, w, h, text, font_size_pt, color_hex,
              bold=False, align="center"):
    """Add a text box with centered text."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = rgb(color_hex)
    run.font.bold = bold
    # Remove internal margin
    from pptx.util import Emu
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    return txBox


def add_line(slide, x1, y1, x2, y2, color_hex, width_pt=1.0, dashed=False):
    """Add a line between two points (in inches)."""
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from lxml import etree

    # python-pptx LINE connector
    x = min(x1, x2)
    y = min(y1, y2)
    w = abs(x2 - x1)
    h = abs(y2 - y1)

    # Degenerate lines (zero width or height) need at least Emu(1)
    w_emu = Inches(w) if w > 0 else Inches(0.001)
    h_emu = Inches(h) if h > 0 else Inches(0.001)

    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = rgb(color_hex)
    connector.line.width = Pt(width_pt)
    if dashed:
        connector.line.dash_style = 4  # DASH
    return connector


# ─────────────────────────────────────────────
# Render Config
# ─────────────────────────────────────────────

@dataclass
class RenderConfig:
    # Slide (LAYOUT_WIDE = 13.3" × 7.5")
    slide_w: float = 13.3
    slide_h: float = 7.5
    bg_color: str  = "F5F6F7"

    # Room
    room_margin_top:    float = 0.35
    room_margin_bottom: float = 0.20
    room_gap:           float = 0.35
    room_pad_top:       float = 0.38
    room_pad_side:      float = 0.22
    room_pad_bottom:    float = 0.15
    room_label_h:       float = 0.26

    # Cabinet column unit
    cab_col_w:          float = 0.92
    cab_gap:            float = 0.16
    cab_label_h:        float = 0.22
    cab_inner_pad:      float = 0.07

    # Tower
    tower_gap:          float = 0.05

    # Devices
    dev_h:              float = 0.28
    dev_gap:            float = 0.05
    cioc_h:             float = 0.32

    # FOPP
    fopp_w:             float = 0.70
    fopp_h:             float = 0.28
    fopp_gap_above:     float = 0.18

    # Colors
    bg:               str = "F5F6F7"
    room_fill:        str = "FFFFFF"
    room_border:      str = "B0BEC5"
    room_lbl_bg:      str = "37474F"
    room_lbl_txt:     str = "FFFFFF"

    srv_fill:         str = "E8EAF6"
    srv_border:       str = "3949AB"
    srv_lbl_bg:       str = "3949AB"
    srv_lbl_txt:      str = "FFFFFF"

    io_fill:          str = "E8F5E9"
    io_border:        str = "2E7D32"
    io_lbl_bg:        str = "2E7D32"
    io_lbl_txt:       str = "FFFFFF"

    wd_fill:          str = "FFF8E1"
    wd_border:        str = "F57F17"
    wd_lbl_bg:        str = "F57F17"
    wd_lbl_txt:       str = "FFFFFF"

    tower_fill:       str = "F1F8E9"
    tower_border:     str = "81C784"

    cioc_fill:        str = "C8E6C9"
    cioc_border:      str = "388E3C"
    cioc_txt:         str = "1B5E20"

    charm_fill:       str = "DCEDC8"
    charm_border:     str = "AED581"
    charm_txt:        str = "33691E"

    cntr_fill:        str = "E3F2FD"
    cntr_border:      str = "1565C0"
    cntr_txt:         str = "0D47A1"

    enc_fill:         str = "F3E5F5"
    enc_border:       str = "7B1FA2"
    enc_txt:          str = "4A148C"

    fopp_fill:        str = "FFF3E0"
    fopp_border:      str = "E65100"
    fopp_txt:         str = "BF360C"

    conn_intra:       str = "78909C"
    conn_inter:       str = "E65100"


# ─────────────────────────────────────────────
# Coordinate Engine
# ─────────────────────────────────────────────

class CoordinateEngine:
    """
    Converts Layout JSON grid positions → absolute x/y/w/h in inches.
    Returns pos: dict[id → {x, y, w, h}]
    """

    def __init__(self, cfg: RenderConfig):
        self.cfg = cfg

    def build(self, elements: list) -> dict:
        cfg = self.cfg
        pos = {}

        # Group by room
        rooms: Dict[int, dict] = {}
        for el in elements:
            ri = el["room_index"]
            if ri not in rooms:
                rooms[ri] = {"cabinets": [], "fopps": [], "devices": []}
            et = el["element_type"]
            if et == "CABINET":   rooms[ri]["cabinets"].append(el)
            elif et == "FOPP":    rooms[ri]["fopps"].append(el)
            elif et == "DEVICE":  rooms[ri]["devices"].append(el)

        # Room height (same for all)
        room_h = cfg.slide_h - cfg.room_margin_top - cfg.room_margin_bottom
        room_y = cfg.room_margin_top

        # Cabinet height: fills available room space
        cab_available_h = (room_h
                           - cfg.room_pad_top
                           - cfg.room_pad_bottom
                           - cfg.fopp_h
                           - cfg.fopp_gap_above)
        cab_y = room_y + cfg.room_pad_top

        # Room widths
        room_widths = {}
        for ri, room in rooms.items():
            total_cols = max((c["col"] + c["col_span"] for c in room["cabinets"]), default=1)
            room_widths[ri] = (cfg.room_pad_side * 2
                               + total_cols * cfg.cab_col_w
                               + max(total_cols - 1, 0) * cfg.cab_gap)

        # Room x offsets (centred on slide)
        n_rooms = len(rooms)
        total_w = sum(room_widths.values()) + (n_rooms - 1) * cfg.room_gap
        room_x_off = {}
        cur_x = (cfg.slide_w - total_w) / 2
        for ri in sorted(rooms.keys()):
            room_x_off[ri] = cur_x
            pos[f"__room_{ri}"] = {
                "x": cur_x, "y": room_y,
                "w": room_widths[ri], "h": room_h,
                "room_index": ri
            }
            cur_x += room_widths[ri] + cfg.room_gap

        # Cabinets and devices
        for ri, room in rooms.items():
            rx = room_x_off[ri]

            for cab in room["cabinets"]:
                cab_x = rx + cfg.room_pad_side + cab["col"] * (cfg.cab_col_w + cfg.cab_gap)
                cab_w = (cab["col_span"] * cfg.cab_col_w
                         + (cab["col_span"] - 1) * cfg.cab_gap)
                pos[cab["id"]] = {
                    "x": cab_x, "y": cab_y,
                    "w": cab_w,  "h": cab_available_h
                }

                # Devices inside cabinet
                devs = [e for e in elements
                        if e["element_type"] == "DEVICE" and e.get("parent") == cab["id"]]

                if cab["cabinet_type"] == "IO":
                    for dev in devs:
                        tc = dev["tower_col"]
                        tower_x = cab_x + tc * (cfg.cab_col_w + cfg.tower_gap)
                        dev_base_y = cab_y + cfg.cab_label_h + cfg.cab_inner_pad

                        if dev["device_type"] == "CIOC":
                            pos[dev["id"]] = {
                                "x": tower_x + cfg.cab_inner_pad,
                                "y": dev_base_y,
                                "w": cfg.cab_col_w - cfg.cab_inner_pad * 2,
                                "h": cfg.cioc_h
                            }
                        else:
                            # CHARM: stack_index 1..N (0 = CIOC)
                            bp_y = (dev_base_y + cfg.cioc_h + cfg.dev_gap
                                    + (dev["stack_index"] - 1) * (cfg.dev_h + cfg.dev_gap))
                            pos[dev["id"]] = {
                                "x": tower_x + cfg.cab_inner_pad,
                                "y": bp_y,
                                "w": cfg.cab_col_w - cfg.cab_inner_pad * 2,
                                "h": cfg.dev_h
                            }
                else:
                    # SERVER / WORKDESK — flat stack
                    for dev in devs:
                        dev_y = (cab_y + cfg.cab_label_h + cfg.cab_inner_pad
                                 + dev["stack_index"] * (cfg.dev_h + cfg.dev_gap))
                        pos[dev["id"]] = {
                            "x": cab_x + cfg.cab_inner_pad,
                            "y": dev_y,
                            "w": cab_w - cfg.cab_inner_pad * 2,
                            "h": cfg.dev_h
                        }

            # FOPP — bottom-right of room
            for fopp in room["fopps"]:
                fopp_y = cab_y + cab_available_h + cfg.fopp_gap_above
                fopp_x = rx + room_widths[ri] - cfg.room_pad_side - cfg.fopp_w
                pos[fopp["id"]] = {
                    "x": fopp_x, "y": fopp_y,
                    "w": cfg.fopp_w, "h": cfg.fopp_h
                }

        return pos


# ─────────────────────────────────────────────
# PPT Renderer
# ─────────────────────────────────────────────

class PptRenderer:

    ROOM_LABELS = {
        "OTHER":    "PDC ROOM",
        "OPERATOR": "OPERATOR ROOM",
    }

    CAB_COLORS = {
        "SERVER":   ("srv_fill",  "srv_border",  "srv_lbl_bg",  "srv_lbl_txt"),
        "IO":       ("io_fill",   "io_border",   "io_lbl_bg",   "io_lbl_txt"),
        "WORKDESK": ("wd_fill",   "wd_border",   "wd_lbl_bg",   "wd_lbl_txt"),
    }

    DEV_COLORS = {
        "CNTR":     ("cntr_fill",  "cntr_border",  "cntr_txt"),
        "CIOC":     ("cioc_fill",  "cioc_border",  "cioc_txt"),
        "CHARM":    ("charm_fill", "charm_border", "charm_txt"),
        "ENC_COMP": ("enc_fill",   "enc_border",   "enc_txt"),
    }

    def __init__(self, config: RenderConfig = None):
        self.cfg = config or RenderConfig()

    def render(self, layout_json: dict, output_path: str):
        cfg = self.cfg

        # Build presentation
        prs = Presentation()
        prs.slide_width  = Inches(cfg.slide_w)
        prs.slide_height = Inches(cfg.slide_h)

        blank_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = rgb(cfg.bg)

        elements = layout_json["elements"]
        connections = layout_json["connections"]

        # Build coordinate map
        coord_eng = CoordinateEngine(cfg)
        pos = coord_eng.build(elements)

        # Index elements by id
        by_id = {e["id"]: e for e in elements}

        # ── Draw order: rooms → cabinets → towers → devices → FOPPs → connections

        self._draw_rooms(slide, elements, pos, cfg)
        self._draw_cabinets(slide, elements, pos, cfg)
        self._draw_devices(slide, elements, pos, cfg)
        self._draw_fopps(slide, elements, pos, cfg)
        self._draw_connections(slide, connections, pos, by_id, cfg)

        prs.save(output_path)
        print(f"Saved: {output_path}")

    # ── Rooms ──────────────────────────────────

    def _draw_rooms(self, slide, elements, pos, cfg):
        drawn = set()
        for el in elements:
            ri = el["room_index"]
            key = f"__room_{ri}"
            if key in drawn or key not in pos:
                continue
            drawn.add(key)
            p = pos[key]
            room_type = el.get("room_type", "OTHER")

            # Room box
            add_rect(slide, p["x"], p["y"], p["w"], p["h"],
                     cfg.room_fill, cfg.room_border, border_width_pt=1.5)

            # Label bar
            add_rect(slide, p["x"], p["y"], p["w"], cfg.room_label_h,
                     cfg.room_lbl_bg, cfg.room_lbl_bg, border_width_pt=0)
            add_label(slide, p["x"], p["y"], p["w"], cfg.room_label_h,
                      self.ROOM_LABELS.get(room_type, f"ROOM {ri}"),
                      font_size_pt=9, color_hex=cfg.room_lbl_txt, bold=True)

    # ── Cabinets ────────────────────────────────

    def _draw_cabinets(self, slide, elements, pos, cfg):
        cabinets = [e for e in elements if e["element_type"] == "CABINET"]
        all_devices = [e for e in elements if e["element_type"] == "DEVICE"]

        for cab in cabinets:
            p = pos[cab["id"]]
            ct = cab["cabinet_type"]
            fc, bc, lbg, ltxt = (getattr(cfg, k) for k in self.CAB_COLORS.get(ct, self.CAB_COLORS["SERVER"]))

            # Cabinet body
            add_rect(slide, p["x"], p["y"], p["w"], p["h"], fc, bc, border_width_pt=1.2)

            # Label bar
            add_rect(slide, p["x"], p["y"], p["w"], cfg.cab_label_h,
                     lbg, lbg, border_width_pt=0)
            add_label(slide, p["x"], p["y"], p["w"], cfg.cab_label_h,
                      cab["id"], font_size_pt=7, color_hex=ltxt, bold=True)

            # IO: draw tower sub-columns
            if ct == "IO":
                tower_cols = sorted(set(
                    d["tower_col"] for d in all_devices
                    if d.get("parent") == cab["id"]
                ))
                for tc in tower_cols:
                    tx = p["x"] + tc * (cfg.cab_col_w + cfg.tower_gap)
                    ty = p["y"] + cfg.cab_label_h
                    th = p["h"] - cfg.cab_label_h
                    add_rect(slide, tx, ty, cfg.cab_col_w, th,
                             cfg.tower_fill, cfg.tower_border, border_width_pt=0.7)

    # ── Devices ─────────────────────────────────

    def _draw_devices(self, slide, elements, pos, cfg):
        devices = [e for e in elements if e["element_type"] == "DEVICE"]

        for dev in devices:
            p = pos.get(dev["id"])
            if not p:
                continue
            dt = dev["device_type"]
            fc, bc, tc = (getattr(cfg, k) for k in self.DEV_COLORS.get(dt, self.DEV_COLORS["ENC_COMP"]))

            add_rect(slide, p["x"], p["y"], p["w"], p["h"], fc, bc, border_width_pt=0.6)
            add_label(slide, p["x"], p["y"], p["w"], p["h"],
                      self._dev_label(dev), font_size_pt=5.5, color_hex=tc)

    def _dev_label(self, dev):
        dt = dev["device_type"]
        did = dev["id"]
        if dt == "CNTR":
            return did.replace("CNTR-", "") + " CNTR"
        if dt == "CIOC":
            return did
        if dt == "CHARM":
            return did.replace("CHARM-", "")
        if dt == "ENC_COMP":
            # Trim slug: take part after first dash segment, replace underscores
            parts = did.split("-")
            label = " ".join(parts[1:3]).replace("_", " ")
            return label[:16]
        return did

    # ── FOPP nodes ──────────────────────────────

    def _draw_fopps(self, slide, elements, pos, cfg):
        fopps = [e for e in elements if e["element_type"] == "FOPP"]
        for fopp in fopps:
            p = pos[fopp["id"]]
            add_rect(slide, p["x"], p["y"], p["w"], p["h"],
                     cfg.fopp_fill, cfg.fopp_border, border_width_pt=1.5)
            add_label(slide, p["x"], p["y"], p["w"], p["h"],
                      fopp["id"], font_size_pt=6.5, color_hex=cfg.fopp_txt, bold=True)

    # ── Connections ─────────────────────────────

    def _draw_connections(self, slide, connections, pos, by_id, cfg):
        for conn in connections:
            fp = pos.get(conn["from"])
            tp = pos.get(conn["to"])
            if not fp or not tp:
                continue

            if conn["link_type"] == "INTER_ROOM_FIBER":
                self._draw_inter_room(slide, conn, fp, tp, cfg)
            elif conn["link_type"] == "FOPP_TO_CABINET":
                self._draw_fopp_to_cabinet(slide, conn, fp, tp, cfg)

    def _cx(self, p): return p["x"] + p["w"] / 2
    def _cy(self, p): return p["y"] + p["h"] / 2
    def _bottom(self, p): return p["y"] + p["h"]
    def _right(self, p):  return p["x"] + p["w"]

    def _draw_inter_room(self, slide, conn, fp, tp, cfg):
        """Horizontal dashed line between two FOPP nodes in different rooms."""
        # right-center of FROM → left-center of TO
        x1 = self._right(fp)
        y1 = self._cy(fp)
        x2 = tp["x"]
        y2 = self._cy(tp)
        line_y = (y1 + y2) / 2

        # Horizontal segment
        add_line(slide, x1, line_y, x2, line_y,
                 cfg.conn_inter, width_pt=2.0, dashed=True)
        # Verticals to connect each FOPP center to line_y if they differ
        if abs(y1 - line_y) > 0.01:
            add_line(slide, x1, y1, x1, line_y, cfg.conn_inter, width_pt=2.0)
        if abs(y2 - line_y) > 0.01:
            add_line(slide, x2, y2, x2, line_y, cfg.conn_inter, width_pt=2.0)

        # Label
        mid_x = (x1 + x2) / 2
        add_label(slide, mid_x - 0.5, line_y - 0.18, 1.0, 0.16,
                  "INTER-ROOM FIBER", font_size_pt=5.5,
                  color_hex=cfg.conn_inter, bold=True)

    def _draw_fopp_to_cabinet(self, slide, conn, fp, tp, cfg):
        """
        Elbow connector: FOPP top → mid Y → cabinet bottom-center.
        Route: vertical from FOPP center-top up to mid_y,
               horizontal to cabinet center,
               vertical up to cabinet bottom.
        """
        fopp_cx = self._cx(fp)
        fopp_top = fp["y"]
        cab_cx   = self._cx(tp)
        cab_bot  = self._bottom(tp)
        mid_y    = fopp_top - cfg.fopp_gap_above * 0.5

        # FOPP top → mid_y
        add_line(slide, fopp_cx, fopp_top, fopp_cx, mid_y,
                 cfg.conn_intra, width_pt=0.9)
        # Horizontal at mid_y
        x_left  = min(fopp_cx, cab_cx)
        x_right = max(fopp_cx, cab_cx)
        if abs(x_right - x_left) > 0.01:
            add_line(slide, x_left, mid_y, x_right, mid_y,
                     cfg.conn_intra, width_pt=0.9)
        # Cabinet bottom → mid_y
        add_line(slide, cab_cx, cab_bot, cab_cx, mid_y,
                 cfg.conn_intra, width_pt=0.9)


# ─────────────────────────────────────────────
# Test harness
# ─────────────────────────────────────────────
#
# LAYOUT_JSON = {
#   "elements": [
#     {"id":"SRV-01","element_type":"CABINET","cabinet_type":"SERVER","room_type":"OTHER","room_index":0,"row":0,"col":0,"col_span":1},
#     {"id":"CNTR-PRI","element_type":"DEVICE","device_type":"CNTR","room_index":0,"parent":"SRV-01","stack_index":0},
#     {"id":"CNTR-SEC","element_type":"DEVICE","device_type":"CNTR","room_index":0,"parent":"SRV-01","stack_index":1},
#     {"id":"ENC-DELTAV_RACK_WORKSTATION-01","element_type":"DEVICE","device_type":"ENC_COMP","room_index":0,"parent":"SRV-01","stack_index":2},
#     {"id":"ENC-DELTAV_SMART_SWITCH-01","element_type":"DEVICE","device_type":"ENC_COMP","room_index":0,"parent":"SRV-01","stack_index":3},
#     {"id":"ENC-19_INCH_SLIDING_SCREEN-01","element_type":"DEVICE","device_type":"ENC_COMP","room_index":0,"parent":"SRV-01","stack_index":4},
#     {"id":"ENC-APC_UPS-01","element_type":"DEVICE","device_type":"ENC_COMP","room_index":0,"parent":"SRV-01","stack_index":5},
#     {"id":"ENC-HIRSCHMANN_INDUSTRIAL_SWITCH-01","element_type":"DEVICE","device_type":"ENC_COMP","room_index":0,"parent":"SRV-01","stack_index":6},
#     {"id":"IO-01","element_type":"CABINET","cabinet_type":"IO","room_type":"OTHER","room_index":0,"row":0,"col":1,"col_span":3},
#     {"id":"CIOC-01","element_type":"DEVICE","device_type":"CIOC","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":0},
#     {"id":"CHARM-BP-01","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":1},
#     {"id":"CHARM-BP-02","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":2},
#     {"id":"CHARM-BP-03","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":3},
#     {"id":"CHARM-BP-04","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":4},
#     {"id":"CHARM-BP-05","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":5},
#     {"id":"CHARM-BP-06","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":6},
#     {"id":"CHARM-BP-07","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":7},
#     {"id":"CHARM-BP-08","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":0,"stack_index":8},
#     {"id":"CIOC-02","element_type":"DEVICE","device_type":"CIOC","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":0},
#     {"id":"CHARM-BP-09","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":1},
#     {"id":"CHARM-BP-10","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":2},
#     {"id":"CHARM-BP-11","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":3},
#     {"id":"CHARM-BP-12","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":4},
#     {"id":"CHARM-BP-13","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":5},
#     {"id":"CHARM-BP-14","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":6},
#     {"id":"CHARM-BP-15","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":7},
#     {"id":"CHARM-BP-16","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":1,"stack_index":8},
#     {"id":"CIOC-03","element_type":"DEVICE","device_type":"CIOC","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":0},
#     {"id":"CHARM-BP-17","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":1},
#     {"id":"CHARM-BP-18","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":2},
#     {"id":"CHARM-BP-19","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":3},
#     {"id":"CHARM-BP-20","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":4},
#     {"id":"CHARM-BP-21","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":5},
#     {"id":"CHARM-BP-22","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":6},
#     {"id":"CHARM-BP-23","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":7},
#     {"id":"CHARM-BP-24","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-01","tower_col":2,"stack_index":8},
#     {"id":"IO-02","element_type":"CABINET","cabinet_type":"IO","room_type":"OTHER","room_index":0,"row":0,"col":4,"col_span":1},
#     {"id":"CIOC-04","element_type":"DEVICE","device_type":"CIOC","room_index":0,"parent":"IO-02","tower_col":0,"stack_index":0},
#     {"id":"CHARM-BP-25","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-02","tower_col":0,"stack_index":1},
#     {"id":"CHARM-BP-26","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-02","tower_col":0,"stack_index":2},
#     {"id":"CHARM-BP-27","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-02","tower_col":0,"stack_index":3},
#     {"id":"CHARM-BP-28","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-02","tower_col":0,"stack_index":4},
#     {"id":"CHARM-BP-29","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-02","tower_col":0,"stack_index":5},
#     {"id":"CHARM-BP-30","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-02","tower_col":0,"stack_index":6},
#     {"id":"CHARM-BP-31","element_type":"DEVICE","device_type":"CHARM","room_index":0,"parent":"IO-02","tower_col":0,"stack_index":7},
#     {"id":"FOPP-OTHER-01","element_type":"FOPP","device_type":"FOPP","room_type":"OTHER","room_index":0,"row":1,"col":4},
#     {"id":"WD-01","element_type":"CABINET","cabinet_type":"WORKDESK","room_type":"OPERATOR","room_index":1,"row":0,"col":0,"col_span":1},
#     {"id":"WD-DELTAV_FULL-SIZED_TOWER_WORKSTATION-01","element_type":"DEVICE","device_type":"ENC_COMP","room_index":1,"parent":"WD-01","stack_index":0},
#     {"id":"WD-24-INCH_MONITOR-01","element_type":"DEVICE","device_type":"ENC_COMP","room_index":1,"parent":"WD-01","stack_index":1},
#     {"id":"WD-PRINTER-01","element_type":"DEVICE","device_type":"ENC_COMP","room_index":1,"parent":"WD-01","stack_index":2},
#     {"id":"FOPP-OPR-01","element_type":"FOPP","device_type":"FOPP","room_type":"OPERATOR","room_index":1,"row":1,"col":0}
#   ],
#   "connections": [
#     {"from":"FOPP-OTHER-01","to":"SRV-01","link_type":"FOPP_TO_CABINET","scope":"INTRA_ROOM"},
#     {"from":"FOPP-OTHER-01","to":"IO-01","link_type":"FOPP_TO_CABINET","scope":"INTRA_ROOM"},
#     {"from":"FOPP-OTHER-01","to":"IO-02","link_type":"FOPP_TO_CABINET","scope":"INTRA_ROOM"},
#     {"from":"FOPP-OPR-01","to":"WD-01","link_type":"FOPP_TO_CABINET","scope":"INTRA_ROOM"},
#     {"from":"FOPP-OTHER-01","to":"FOPP-OPR-01","link_type":"INTER_ROOM_FIBER","scope":"INTER_ROOM"}
#   ]
# }

# if __name__ == "__main__":
#     renderer = PptRenderer()
#     renderer.render(LAYOUT_JSON, "deltav_architecture.pptx")